from http.server import BaseHTTPRequestHandler
import json
import os
import hmac
import tempfile
import traceback
import shutil
import urllib.request
import urllib.parse
import zipfile
from copy import deepcopy

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload

# PowerPoint XML namespaces
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
P159_NS = 'http://schemas.microsoft.com/office/powerpoint/2015/09/main'
MC_NS = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _pn(tag):
    """Build a namespaced tag for presentationml namespace."""
    return f'{{{P_NS}}}{tag}'


def _p14n(tag):
    """Build a namespaced tag for PowerPoint 2010 namespace."""
    return f'{{{P14_NS}}}{tag}'


def _rn(tag):
    """Build a namespaced tag for relationships namespace."""
    return f'{{{R_NS}}}{tag}'


def verify_auth(request):
    """Verify Bearer token using timing-safe comparison."""
    auth_header = request.headers.get('Authorization', '')
    expected = f"Bearer {os.environ.get('AUTH_SECRET', '')}"
    if not auth_header or not hmac.compare_digest(auth_header, expected):
        return False
    return True


def get_drive_service():
    creds_json = json.loads(os.environ['GOOGLE_SERVICE_ACCOUNT_JSON'])
    credentials = service_account.Credentials.from_service_account_info(
        creds_json,
        scopes=['https://www.googleapis.com/auth/drive']
    )
    return build('drive', 'v3', credentials=credentials)


def download_file_by_id(service, file_id, dest_path):
    """Download a file from Google Drive by its file ID."""
    request = service.files().get_media(fileId=file_id, supportsAllDrives=True)
    with open(dest_path, 'wb') as f:
        downloader = MediaIoBaseDownload(f, request)
        done = False
        while not done:
            _, done = downloader.next_chunk()


def overwrite_drive_file(service, file_id, file_path):
    """Overwrite an existing file on Google Drive (update in-place)."""
    media = MediaFileUpload(
        file_path,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    file = service.files().update(
        fileId=file_id,
        media_body=media,
        fields='id, name, webViewLink',
        supportsAllDrives=True,
    ).execute()
    return file


def upload_to_blob(file_path, file_name):
    """Upload a file to Vercel Blob API and return the URL.

    Args:
        file_path: Path to the file to upload
        file_name: Desired filename in Blob storage

    Returns:
        str: The URL of the uploaded blob

    Raises:
        ValueError: If BLOB_READ_WRITE_TOKEN is not set
        Exception: If upload fails
    """
    token = os.environ.get('BLOB_READ_WRITE_TOKEN')
    if not token:
        raise ValueError('BLOB_READ_WRITE_TOKEN environment variable is not set')

    # Read file bytes
    with open(file_path, 'rb') as f:
        file_bytes = f.read()

    # Construct Blob API URL
    encoded_name = urllib.parse.quote(file_name, safe='')
    blob_url = f'https://blob.vercel-storage.com/pptx-exports/{encoded_name}'

    # Create request
    req = urllib.request.Request(
        blob_url,
        data=file_bytes,
        method='PUT',
        headers={
            'authorization': f'Bearer {token}',
            'x-api-version': '7',
            'x-content-type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        }
    )

    # Upload
    try:
        with urllib.request.urlopen(req) as response:
            response_data = json.loads(response.read().decode('utf-8'))
            return response_data.get('url', '')
    except urllib.error.HTTPError as e:
        error_body = e.read().decode('utf-8') if e.fp else 'No error body'
        raise Exception(f'Blob upload failed with status {e.code}: {error_body}')
    except Exception as e:
        raise Exception(f'Blob upload failed: {str(e)}')


def cleanup_orphaned_parts(pptx_path):
    """Remove orphaned parts and broken references from PPTX ZIP.

    python-pptx does not fully clean up deleted slides during serialization.
    This post-processing step:
    1. Removes slide XML files not referenced in presentation.xml.rels
    2. Removes notesSlide XML files not referenced by any remaining slide
    3. Strips broken notesSlide references from slide .rels files
    4. Updates [Content_Types].xml to match
    """
    from lxml import etree as _et

    tmp_path = pptx_path + '.tmp'
    ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'
    rels_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'

    with zipfile.ZipFile(pptx_path, 'r') as zin:
        names = set(zin.namelist())

        # Step 1: Find slides referenced by presentation.xml.rels
        prs_rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
        prs_rels_root = _et.fromstring(prs_rels_xml)
        prs_referenced = set()
        for rel in prs_rels_root.findall(f'{{{rels_ns}}}Relationship'):
            target = rel.get('Target', '')
            if target.startswith('/'):
                prs_referenced.add(target.lstrip('/'))
            else:
                prs_referenced.add('ppt/' + target)

        # Step 2: Find orphaned slide files
        orphaned = set()
        for name in names:
            if (name.startswith('ppt/slides/slide') and
                    name.endswith('.xml') and '/_rels/' not in name
                    and name not in prs_referenced):
                orphaned.add(name)

        # Step 3: Collect notesSlide files actually in the ZIP
        notes_in_zip = set()
        for name in names:
            if (name.startswith('ppt/notesSlides/notesSlide') and
                    name.endswith('.xml') and '/_rels/' not in name):
                notes_in_zip.add(name)

        # Step 4: Find notesSlides referenced by remaining (non-orphaned) slides
        notes_referenced = set()
        for name in names:
            if (name.startswith('ppt/slides/_rels/slide') and
                    name.endswith('.xml.rels')):
                # Skip rels for orphaned slides
                slide_name = name.replace('/_rels/', '/').removesuffix('.rels')
                if slide_name in orphaned:
                    continue
                rels_xml = zin.read(name)
                rels_root = _et.fromstring(rels_xml)
                for rel in rels_root.findall(f'{{{rels_ns}}}Relationship'):
                    target = rel.get('Target', '')
                    if 'notesSlide' in target:
                        resolved = os.path.normpath(
                            os.path.join('ppt/slides', target)
                        ).replace('\\', '/')
                        notes_referenced.add(resolved)

        # Orphaned notesSlides = in ZIP but not referenced by any remaining slide
        orphaned_notes = notes_in_zip - notes_referenced
        orphaned |= orphaned_notes

        # Build set of all files to remove (orphaned parts + their .rels)
        to_remove = set(orphaned)
        for orph in orphaned:
            d, f = orph.rsplit('/', 1)
            to_remove.add(d + '/_rels/' + f + '.rels')

        # Also find notesSlide files referenced but missing from ZIP
        # (broken references that need stripping from slide .rels)
        missing_notes = notes_referenced - notes_in_zip

        if not to_remove and not missing_notes:
            return  # Nothing to clean

        # Rewrite ZIP
        with zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename in to_remove:
                    continue

                data = zin.read(item.filename)

                # Strip broken notesSlide refs from slide .rels files
                if (missing_notes and
                        item.filename.startswith('ppt/slides/_rels/') and
                        item.filename.endswith('.xml.rels')):
                    rels_root = _et.fromstring(data)
                    changed = False
                    for rel in rels_root.findall(f'{{{rels_ns}}}Relationship'):
                        target = rel.get('Target', '')
                        if 'notesSlide' in target:
                            resolved = os.path.normpath(
                                os.path.join('ppt/slides', target)
                            ).replace('\\', '/')
                            if resolved in missing_notes or resolved in orphaned_notes:
                                rels_root.remove(rel)
                                changed = True
                    if changed:
                        data = _et.tostring(rels_root, xml_declaration=True,
                                            encoding='UTF-8', standalone=True)

                # Remove orphaned entries from Content_Types
                if item.filename == '[Content_Types].xml':
                    ct_root = _et.fromstring(data)
                    for override in list(ct_root.findall(f'{{{ct_ns}}}Override')):
                        pn = override.get('PartName', '').lstrip('/')
                        if pn in orphaned:
                            ct_root.remove(override)
                    data = _et.tostring(ct_root, xml_declaration=True,
                                        encoding='UTF-8', standalone=True)

                zout.writestr(item, data)

    os.replace(tmp_path, pptx_path)


def parse_sections(prs):
    """Parse sections from presentation.xml.

    Supports both standard <p:sectionLst> and PowerPoint 2010's
    <p14:sectionLst> inside <p:extLst>.
    """
    prs_xml = prs._element

    # Try standard namespace first
    section_lst = prs_xml.find(_pn('sectionLst'))
    ns_fn = _pn

    # Fall back to p14 namespace in extLst
    if section_lst is None:
        ext_lst = prs_xml.find(_pn('extLst'))
        if ext_lst is not None:
            for ext in ext_lst.findall(_pn('ext')):
                section_lst = ext.find(_p14n('sectionLst'))
                if section_lst is not None:
                    ns_fn = _p14n
                    break

    if section_lst is None:
        raise ValueError("Template has no sections. The template must use PowerPoint sections.")

    sections = []
    for section_el in section_lst.findall(ns_fn('section')):
        name = section_el.get('name', '')
        sect_id = section_el.get('id', '')
        sld_id_lst = section_el.find(ns_fn('sldIdLst'))
        slide_ids = []
        if sld_id_lst is not None:
            for sld_id_el in sld_id_lst.findall(ns_fn('sldId')):
                slide_ids.append(int(sld_id_el.get('id')))
        sections.append({
            'name': name,
            'id': sect_id,
            'slide_ids': slide_ids,
            'element': section_el,
            'ns_fn': ns_fn,
        })

    return sections


def find_section_by_name(sections, name):
    """Find a section dict by exact name match."""
    for s in sections:
        if s['name'] == name:
            return s
    return None


def get_slide_id_map(prs):
    """Build a mapping of slide_id (int) -> (slide_index, rId, slide object)."""
    prs_xml = prs._element
    sld_id_lst = prs_xml.find(_pn('sldIdLst'))

    sld_entries = []
    for sld_id_el in sld_id_lst.findall(_pn('sldId')):
        sid = int(sld_id_el.get('id'))
        rid = sld_id_el.get(_rn('id'))
        sld_entries.append((sid, rid, sld_id_el))

    slide_map = {}
    for idx, (sid, rid, el) in enumerate(sld_entries):
        slide_map[sid] = {
            'index': idx,
            'rId': rid,
            'element': el,
            'slide': prs.slides[idx] if idx < len(prs.slides) else None,
        }

    return slide_map


def _remap_slide_rids(element, rId_map):
    """Remap relationship ID references in slide XML after cloning.

    Scans all attributes in the relationships namespace (r:id, r:embed,
    r:link, etc.) and replaces old rId values with the new slide's rIds.
    """
    if not rId_map or all(k == v for k, v in rId_map.items()):
        return
    for el in element.iter():
        for attr_name in list(el.attrib.keys()):
            if R_NS in attr_name:
                val = el.get(attr_name)
                if val in rId_map:
                    el.set(attr_name, rId_map[val])


def duplicate_slide(prs, slide):
    """Clone a slide and append it to the presentation.

    Returns: (new_slide, new_slide_id, new_entry_element)
    """
    # Find highest slide part number to avoid ZIP duplicate name collisions.
    # add_slide() uses len(prs.slides)+1 which can collide with existing
    # slide partnames when slides have been deleted from the collection.
    max_num = 0
    for rel in prs.part.rels.values():
        if rel.is_external:
            continue
        pn = str(rel.target_part.partname)
        if pn.startswith('/ppt/slides/slide') and pn.endswith('.xml'):
            try:
                num = int(pn[17:-4])
                max_num = max(max_num, num)
            except ValueError:
                pass

    new_slide = prs.slides.add_slide(slide.slide_layout)

    # Rename part to guaranteed-unique name if add_slide assigned a conflicting one
    safe_name = PackURI('/ppt/slides/slide%d.xml' % (max_num + 1))
    if new_slide.part.partname != safe_name:
        new_slide.part.partname = safe_name

    # Replace new slide's XML content with a copy of the source
    for child in list(new_slide._element):
        new_slide._element.remove(child)
    for child in deepcopy(slide._element):
        new_slide._element.append(child)

    # Clear lazyproperty caches that were populated by add_slide().
    # add_slide() calls clone_layout_placeholders() which accesses .shapes,
    # caching a SlideShapes referencing the original spTree. After replacing
    # _element's children with a deepcopy, these caches point to orphaned XML.
    # Clearing them forces reconstruction from the live _element on next access.
    for attr in ('shapes', 'placeholders', 'background'):
        new_slide.__dict__.pop(attr, None)

    # Build rId mapping (source rId → new rId) and copy relationships.
    # The copied XML references the source slide's rIds, but get_or_add()
    # may assign different rIds on the new slide. We must remap the XML
    # references to match, otherwise PowerPoint repair strips the content.
    rId_map = {}

    # Map the layout relationship (already created by add_slide)
    for src_rId, src_rel in slide.part.rels.items():
        if src_rel.reltype == RT.SLIDE_LAYOUT:
            for new_rId, new_rel in new_slide.part.rels.items():
                if new_rel.reltype == RT.SLIDE_LAYOUT:
                    rId_map[src_rId] = new_rId
                    break
            break

    # Copy non-layout relationships and record the mapping.
    # Skip notesSlide relationships — generated slides don't need speaker notes,
    # and copying them creates broken references when the base slide is deleted.
    for src_rId, src_rel in slide.part.rels.items():
        if src_rel.reltype == RT.SLIDE_LAYOUT:
            continue
        if src_rel.reltype == RT.NOTES_SLIDE:
            continue
        if src_rel.is_external:
            new_rId = new_slide.part.rels.get_or_add_ext_rel(
                src_rel.reltype, src_rel.target_ref
            )
        else:
            new_rId = new_slide.part.rels.get_or_add(
                src_rel.reltype, src_rel.target_part
            )
        rId_map[src_rId] = new_rId

    # Remap all r:* attribute references in the copied XML
    _remap_slide_rids(new_slide._element, rId_map)

    prs_xml = prs._element
    sld_id_lst = prs_xml.find(_pn('sldIdLst'))

    sld_id_entries = sld_id_lst.findall(_pn('sldId'))
    new_entry = sld_id_entries[-1]
    new_slide_id = int(new_entry.get('id'))

    return new_slide, new_slide_id, new_entry


def delete_slide_by_id(prs, slide_id):
    """Remove a slide from the presentation by its numeric slide ID."""
    prs_xml = prs._element
    sld_id_lst = prs_xml.find(_pn('sldIdLst'))

    target_el = None
    target_rId = None
    for sld_id_el in sld_id_lst.findall(_pn('sldId')):
        if int(sld_id_el.get('id')) == slide_id:
            target_el = sld_id_el
            target_rId = sld_id_el.get(_rn('id'))
            break

    if target_el is None:
        raise ValueError(f"Slide with id={slide_id} not found in presentation")

    prs_part = prs.part
    # Clear the deleted slide's own relationships to prevent ghost parts
    slide_part = prs_part.rels[target_rId].target_part
    for key in list(slide_part.rels.keys()):
        slide_part.rels.pop(key)

    prs_part.rels.pop(target_rId)

    sld_id_lst.remove(target_el)


def move_slide_id_after(prs, slide_id_to_move, after_slide_id):
    """Move a <p:sldId> entry to be positioned right after another slide ID."""
    prs_xml = prs._element
    sld_id_lst = prs_xml.find(_pn('sldIdLst'))

    move_el = None
    after_el = None
    for sld_id_el in sld_id_lst.findall(_pn('sldId')):
        sid = int(sld_id_el.get('id'))
        if sid == slide_id_to_move:
            move_el = sld_id_el
        if sid == after_slide_id:
            after_el = sld_id_el

    if move_el is None or after_el is None:
        raise ValueError(f"Could not find slide IDs for reordering: move={slide_id_to_move}, after={after_slide_id}")

    sld_id_lst.remove(move_el)

    after_index = list(sld_id_lst).index(after_el)
    sld_id_lst.insert(after_index + 1, move_el)


def inject_text_into_shape(shape, text):
    """Replace text in a shape while preserving formatting."""
    tf = shape.text_frame

    template_pPr = None
    template_rPr = None
    if tf.paragraphs:
        first_para = tf.paragraphs[0]
        pPr = first_para._p.find(qn('a:pPr'))
        if pPr is not None:
            template_pPr = deepcopy(pPr)
        if first_para.runs:
            rPr = first_para.runs[0]._r.find(qn('a:rPr'))
            if rPr is not None:
                template_rPr = deepcopy(rPr)

    p_elements = tf._txBody.findall(qn('a:p'))
    for p in p_elements[1:]:
        tf._txBody.remove(p)

    lines = text.split('\n') if text else ['']

    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
            for r in para._p.findall(qn('a:r')):
                para._p.remove(r)
        else:
            para = tf.add_paragraph()

        if template_pPr is not None:
            existing_pPr = para._p.find(qn('a:pPr'))
            if existing_pPr is not None:
                para._p.remove(existing_pPr)
            para._p.insert(0, deepcopy(template_pPr))

        run = para.add_run()
        run.text = line

        if template_rPr is not None:
            existing_rPr = run._r.find(qn('a:rPr'))
            if existing_rPr is not None:
                run._r.remove(existing_rPr)
            run._r.insert(0, deepcopy(template_rPr))


def get_first_textbox(slide):
    """Get the first shape with a text_frame on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            return shape
    return None


def get_largest_textbox(slide):
    """Get the largest text-bearing shape on a slide."""
    textboxes = [shape for shape in slide.shapes if shape.has_text_frame]
    if not textboxes:
        return None

    return max(
        textboxes,
        key=lambda shape: int(shape.width or 0) * int(shape.height or 0)
    )


def clear_other_textboxes(slide, keep_shape):
    """Clear text-bearing shapes except the selected body textbox."""
    keep_element = keep_shape._element if keep_shape is not None else None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape._element is keep_element:
            continue
        inject_text_into_shape(shape, '')


def strip_textbox_numbering(shape):
    """Disable inherited PowerPoint paragraph numbering on a text shape."""
    if not shape.has_text_frame:
        return

    numbering_tags = {
        qn('a:buNone'),
        qn('a:buAutoNum'),
        qn('a:buChar'),
        qn('a:buBlip'),
    }
    insert_before_tags = {
        qn('a:tabLst'),
        qn('a:defRPr'),
        qn('a:extLst'),
    }

    for paragraph in shape.text_frame.paragraphs:
        pPr = paragraph._p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.Element(qn('a:pPr'))
            paragraph._p.insert(0, pPr)

        for child in list(pPr):
            if child.tag in numbering_tags:
                pPr.remove(child)

        bu_none = etree.Element(qn('a:buNone'))
        for index, child in enumerate(list(pPr)):
            if child.tag in insert_before_tags:
                pPr.insert(index, bu_none)
                break
        else:
            pPr.append(bu_none)


def set_slide_notes(slide, text):
    """Add speaker notes to a slide using python-pptx."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def set_morph_transition(slide, duration_ms=1000):
    """Set a morph transition on a slide using mc:AlternateContent.

    Uses the same backward-compatible pattern as PowerPoint itself:
    mc:Choice with p159:morph for 2019+, mc:Fallback with p:fade for older.
    """
    sld = slide._element
    # Remove existing bare transitions and mc:AlternateContent wrappers
    for existing in sld.findall(_pn('transition')):
        sld.remove(existing)
    for existing in sld.findall(f'{{{MC_NS}}}AlternateContent'):
        sld.remove(existing)

    transition_xml = (
        '<mc:AlternateContent'
        f' xmlns:mc="{MC_NS}"'
        f' xmlns:p159="{P159_NS}">'
        '<mc:Choice Requires="p159">'
        f'<p:transition xmlns:p="{P_NS}" spd="slow">'
        '<p159:morph option="byObject"/>'
        '</p:transition>'
        '</mc:Choice>'
        '<mc:Fallback xmlns="">'
        f'<p:transition xmlns:p="{P_NS}" spd="slow">'
        '<p:fade/>'
        '</p:transition>'
        '</mc:Fallback>'
        '</mc:AlternateContent>'
    )
    sld.append(etree.fromstring(transition_xml))


def normalize_reference_text(reference):
    """Normalize reference whitespace for title slide injection."""
    return ' '.join(str(reference or '').split())


def process_scripture_section(prs, scripture, section, slide_id_map):
    """Process scripture: inject reference, clone body base per page, update section."""
    slide_ids = section['slide_ids']

    if len(slide_ids) < 2:
        raise ValueError(
            f"Section '{section['name']}' needs at least 2 slides (title + base), "
            f"but has {len(slide_ids)}"
        )

    title_slide_id = slide_ids[0]
    body_base_slide_id = slide_ids[1]

    title_slide = slide_id_map[title_slide_id]['slide']
    title_shape = get_first_textbox(title_slide)
    if title_shape:
        inject_text_into_shape(
            title_shape,
            normalize_reference_text(scripture.get('reference', ''))
        )

    body_base_slide = slide_id_map[body_base_slide_id]['slide']

    for sid in slide_ids[2:]:
        delete_slide_by_id(prs, sid)

    generated_slide_ids = []
    last_slide_id = body_base_slide_id
    pages = scripture.get('pages', [])

    for page_idx, page in enumerate(pages, 1):
        new_slide, new_sid, new_el = duplicate_slide(prs, body_base_slide)
        textbox = get_largest_textbox(new_slide)
        if textbox:
            clear_other_textboxes(new_slide, textbox)
            inject_text_into_shape(textbox, page.get('text', ''))
            strip_textbox_numbering(textbox)
        move_slide_id_after(prs, new_sid, last_slide_id)
        generated_slide_ids.append(new_sid)
        last_slide_id = new_sid

        note_title = page.get('title') or f"{section['name']}-{page_idx}"
        set_slide_notes(new_slide, note_title)

    delete_slide_by_id(prs, body_base_slide_id)

    section_el = section['element']
    ns_fn = section.get('ns_fn', _pn)
    sld_id_lst = section_el.find(ns_fn('sldIdLst'))

    for child in list(sld_id_lst):
        sld_id_lst.remove(child)

    title_entry = etree.SubElement(sld_id_lst, ns_fn('sldId'))
    title_entry.set('id', str(title_slide_id))

    for gen_sid in generated_slide_ids:
        entry = etree.SubElement(sld_id_lst, ns_fn('sldId'))
        entry.set('id', str(gen_sid))

    return len(generated_slide_ids)


def process_song_section(prs, song, section, slide_id_map, shared_base_slide_id):
    """Process a single song: inject title, clone base slide for lyrics, update section."""
    slide_ids = section['slide_ids']

    if len(slide_ids) < 2:
        raise ValueError(
            f"Section '{section['name']}' needs at least 2 slides (title + base), "
            f"but has {len(slide_ids)}"
        )

    title_slide_id = slide_ids[0]
    section_base_slide_id = slide_ids[1]

    title_slide = slide_id_map[title_slide_id]['slide']
    title_shape = get_first_textbox(title_slide)
    if title_shape:
        inject_text_into_shape(title_shape, song['title'])

    base_slide = slide_id_map[shared_base_slide_id]['slide']

    slides_to_delete = slide_ids[2:]
    for sid in slides_to_delete:
        delete_slide_by_id(prs, sid)

    section_order = song.get('section_order', [])
    lyrics = song.get('lyrics', [])
    section_lyrics_map = song.get('section_lyrics_map', {})

    generated_slide_ids = []
    last_slide_id = section_base_slide_id

    for sect_idx, sect_name in enumerate(section_order):
        sect_idx_str = str(sect_idx)

        if sect_name.strip().lower() == 'intro':
            continue

        mapped_lyrics_indices = section_lyrics_map.get(sect_idx_str,
                                section_lyrics_map.get(sect_idx, []))

        # Per-occurrence page count for note labeling
        occur_page_count = max(len(mapped_lyrics_indices), 1)

        if not mapped_lyrics_indices:
            new_slide, new_sid, new_el = duplicate_slide(prs, base_slide)
            textbox = get_first_textbox(new_slide)
            if textbox:
                inject_text_into_shape(textbox, '')
            set_morph_transition(new_slide)
            move_slide_id_after(prs, new_sid, last_slide_id)
            generated_slide_ids.append(new_sid)
            last_slide_id = new_sid
            # Add section note (single page = plain name)
            set_slide_notes(new_slide, sect_name)
        else:
            for page_num, lyrics_idx in enumerate(mapped_lyrics_indices, 1):
                if lyrics_idx >= len(lyrics):
                    raise ValueError(
                        f"Section '{section['name']}', sectionOrder[{sect_idx}]='{sect_name}': "
                        f"lyrics index {lyrics_idx} out of range (have {len(lyrics)} lyrics pages)"
                    )

                lyrics_text = lyrics[lyrics_idx]

                new_slide, new_sid, new_el = duplicate_slide(prs, base_slide)
                textbox = get_first_textbox(new_slide)
                if textbox:
                    inject_text_into_shape(textbox, lyrics_text)
                move_slide_id_after(prs, new_sid, last_slide_id)
                generated_slide_ids.append(new_sid)
                last_slide_id = new_sid
                # Add section note
                if occur_page_count == 1:
                    set_slide_notes(new_slide, sect_name)
                else:
                    set_slide_notes(new_slide, f"{sect_name}-{page_num}")

    # Append a blank slide at the end of the song with morph transition
    blank_slide, blank_sid, blank_el = duplicate_slide(prs, base_slide)
    textbox = get_first_textbox(blank_slide)
    if textbox:
        inject_text_into_shape(textbox, '')
    set_morph_transition(blank_slide)
    move_slide_id_after(prs, blank_sid, last_slide_id)
    generated_slide_ids.append(blank_sid)

    if section_base_slide_id != shared_base_slide_id:
        delete_slide_by_id(prs, section_base_slide_id)

    section_el = section['element']
    ns_fn = section.get('ns_fn', _pn)
    sld_id_lst = section_el.find(ns_fn('sldIdLst'))

    for child in list(sld_id_lst):
        sld_id_lst.remove(child)

    title_entry = etree.SubElement(sld_id_lst, ns_fn('sldId'))
    title_entry.set('id', str(title_slide_id))

    for gen_sid in generated_slide_ids:
        entry = etree.SubElement(sld_id_lst, ns_fn('sldId'))
        entry.set('id', str(gen_sid))

    return len(generated_slide_ids)


def process_all_songs(prs, songs, sections=None, slide_id_map=None):
    """Process all songs in the presentation."""
    if sections is None:
        sections = parse_sections(prs)
    if slide_id_map is None:
        slide_id_map = get_slide_id_map(prs)

    song_sections = []
    for song in songs:
        section_name = song.get('section_name', '')
        if not section_name:
            raise ValueError(f"Song '{song.get('title', '?')}' has no section_name")

        section = find_section_by_name(sections, section_name)
        if section is None:
            available = [s['name'] for s in sections]
            raise ValueError(
                f"Section '{section_name}' not found in template. "
                f"Available sections: {available}"
            )
        song_sections.append((song, section))

    # Find shared base slide from song sections only, so scripture bases cannot be selected.
    shared_base_slide_id = None
    for _, section in song_sections:
        if len(section['slide_ids']) >= 2:
            candidate_id = section['slide_ids'][1]
            candidate_slide = slide_id_map[candidate_id]['slide']
            textbox = get_first_textbox(candidate_slide)
            if textbox is not None and textbox.text_frame.text.strip():
                shared_base_slide_id = candidate_id
                break

    # If no song section had a text-bearing base, fall back to first song section's slide_ids[1]
    if shared_base_slide_id is None:
        for _, section in song_sections:
            if len(section['slide_ids']) >= 2:
                shared_base_slide_id = section['slide_ids'][1]
                break

    if shared_base_slide_id is None:
        raise ValueError("No section has a base slide (slide_ids[1]) to use as shared base")

    total_slides = 0
    songs_processed = 0

    for song, section in song_sections:
        slides = process_song_section(prs, song, section, slide_id_map, shared_base_slide_id)
        total_slides += slides
        songs_processed += 1

        slide_id_map = get_slide_id_map(prs)

    # Delete the shared base slide now that all songs have been cloned from it
    delete_slide_by_id(prs, shared_base_slide_id)

    return {
        'songs_processed': songs_processed,
        'slides_generated': total_slides,
    }


def has_scripture_payload(scripture):
    """Return True when an optional scripture payload should be processed."""
    return isinstance(scripture, dict) and bool(scripture.get('pages'))


def process_export(prs, songs, scripture=None):
    """Process optional scripture before existing song export."""
    sections = parse_sections(prs)
    slide_id_map = get_slide_id_map(prs)
    result = {}

    if has_scripture_payload(scripture):
        section_name = scripture.get('section_name', '')
        if not section_name:
            raise ValueError("Scripture payload has no section_name")

        section = find_section_by_name(sections, section_name)
        if section is None:
            available = [s['name'] for s in sections]
            raise ValueError(
                f"Section '{section_name}' not found in template. "
                f"Available sections: {available}"
            )

        scripture_slides = process_scripture_section(prs, scripture, section, slide_id_map)
        result['scripture_processed'] = True
        result['scripture_slides_generated'] = scripture_slides
        slide_id_map = get_slide_id_map(prs)

    result.update(process_all_songs(prs, songs, sections, slide_id_map))
    return result


def inspect_template(pptx_path):
    """Return the slide/shape/section structure for debugging."""
    prs = Presentation(pptx_path)

    slides = []
    for i, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "has_text_frame": shape.has_text_frame,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            if shape.has_text_frame:
                shape_info["text_preview"] = shape.text_frame.text[:100]
                shape_info["paragraph_count"] = len(shape.text_frame.paragraphs)
            shapes.append(shape_info)
        slides.append({"slide_index": i, "shapes": shapes})

    sections_info = []
    try:
        sections = parse_sections(prs)
        for s in sections:
            sections_info.append({
                "name": s['name'],
                "id": s['id'],
                "slide_ids": s['slide_ids'],
                "slide_count": len(s['slide_ids']),
            })
    except ValueError:
        sections_info = None

    return {
        "slide_count": len(prs.slides),
        "slides": slides,
        "sections": sections_info,
    }


class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            if not verify_auth(self):
                self.send_json(401, {"success": False, "error": "Unauthorized"})
                return

            content_length = int(self.headers.get('Content-Length', 0))
            body = json.loads(self.rfile.read(content_length))

            action = body.get('action')

            if action == 'export_lyrics':
                self._handle_export_lyrics(body)
            else:
                self.send_json(400, {"success": False, "error": f"Unknown action: {action}"})

        except ValueError as e:
            self.send_json(400, {
                "success": False,
                "error": str(e)
            })
        except FileNotFoundError as e:
            self.send_json(404, {
                "success": False,
                "error": str(e)
            })
        except Exception as e:
            traceback.print_exc()
            self.send_json(500, {
                "success": False,
                "error": f"Internal error: {str(e)}"
            })

    def _handle_export_lyrics(self, body):
        """Handle the export_lyrics action."""
        file_id = body.get('file_id')
        if not file_id:
            self.send_json(400, {"success": False, "error": "file_id is required"})
            return

        overwrite = body.get('overwrite', False)
        output_file_name = body.get('output_file_name', '')
        # Note: output_folder_id is kept for backward compatibility but unused since new files go to Vercel Blob
        output_folder_id = body.get('output_folder_id', os.environ.get('GOOGLE_DRIVE_TEMPLATE_FOLDER_ID'))
        songs = body.get('songs', [])

        if not overwrite and not output_file_name:
            self.send_json(400, {"success": False, "error": "output_file_name is required when not overwriting"})
            return

        if not songs:
            self.send_json(400, {"success": False, "error": "No songs provided"})
            return

        service = get_drive_service()

        tmp_dir = tempfile.mkdtemp()
        template_path = os.path.join(tmp_dir, 'template.pptx')
        output_path = os.path.join(tmp_dir, 'output.pptx')

        try:
            download_file_by_id(service, file_id, template_path)

            prs = Presentation(template_path)
            result_stats = process_export(prs, songs, body.get('scripture'))
            prs.save(output_path)
            cleanup_orphaned_parts(output_path)

            if overwrite:
                result = overwrite_drive_file(service, file_id, output_path)
                response_data = {
                    "file_id": result['id'],
                    "file_name": result['name'],
                    "web_view_link": result.get('webViewLink', ''),
                    "songs_processed": result_stats['songs_processed'],
                    "slides_generated": result_stats['slides_generated'],
                }
                if 'scripture_processed' in result_stats:
                    response_data["scripture_processed"] = result_stats['scripture_processed']
                    response_data["scripture_slides_generated"] = result_stats['scripture_slides_generated']

                self.send_json(200, {
                    "success": True,
                    "data": response_data
                })
            else:
                # Upload to Vercel Blob and return JSON with download URL
                try:
                    download_url = upload_to_blob(output_path, output_file_name)
                    response_data = {
                        "file_id": "",
                        "file_name": output_file_name,
                        "web_view_link": "",
                        "download_url": download_url,
                        "songs_processed": result_stats['songs_processed'],
                        "slides_generated": result_stats['slides_generated'],
                    }
                    if 'scripture_processed' in result_stats:
                        response_data["scripture_processed"] = result_stats['scripture_processed']
                        response_data["scripture_slides_generated"] = result_stats['scripture_slides_generated']

                    self.send_json(200, {
                        "success": True,
                        "data": response_data
                    })
                except ValueError as e:
                    # BLOB_READ_WRITE_TOKEN not set
                    self.send_json(500, {
                        "success": False,
                        "error": f"Blob storage configuration error: {str(e)}"
                    })
                except Exception as e:
                    # Upload failed
                    self.send_json(500, {
                        "success": False,
                        "error": f"Failed to upload to blob storage: {str(e)}"
                    })
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    def do_GET(self):
        """Health check / template inspection endpoint."""
        try:
            if not verify_auth(self):
                self.send_json(401, {"success": False, "error": "Unauthorized"})
                return

            action = self.headers.get('X-Action', 'health')

            if action == 'inspect':
                file_id = self.headers.get('X-File-Id', '')

                if not file_id:
                    self.send_json(400, {"success": False, "error": "X-File-Id header required"})
                    return

                service = get_drive_service()
                tmp_dir = tempfile.mkdtemp()
                template_path = os.path.join(tmp_dir, 'template.pptx')

                try:
                    download_file_by_id(service, file_id, template_path)
                    structure = inspect_template(template_path)
                    self.send_json(200, {"success": True, "data": structure})
                finally:
                    shutil.rmtree(tmp_dir, ignore_errors=True)
            else:
                self.send_json(200, {"success": True, "data": {"status": "ok"}})

        except Exception as e:
            self.send_json(500, {"success": False, "error": str(e)})

    def send_json(self, status_code, data):
        self.send_response(status_code)
        self.send_header('Content-Type', 'application/json')
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode('utf-8'))
